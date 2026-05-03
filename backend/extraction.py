"""
Text extraction from PDF and PPTX files.
Returns a list of {text, page_number} dicts per document.
"""
import re
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
import docx  # python-docx


def extract_pdf(file_path: str) -> list[dict]:
    """Extract text page-by-page from a PDF."""
    pages = []
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        if text.strip():
            pages.append({"text": text, "page_number": page_num})
    doc.close()
    return pages


def extract_pptx(file_path: str) -> list[dict]:
    """Extract text slide-by-slide from a PPTX."""
    pages = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            pages.append({"text": "\n".join(texts), "page_number": slide_num})
    return pages


def extract_docx(file_path: str) -> list[dict]:
    """Extract text paragraph-by-paragraph from a DOCX."""
    pages = []
    doc = docx.Document(file_path)
    # Treat the whole document as page 1 for now, or split by some delimiter if needed
    text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    if text:
        pages.append({"text": text, "page_number": 1})
    return pages


def extract_txt(file_path: str) -> list[dict]:
    """Extract text from a plain TXT file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    if text:
        return [{"text": text, "page_number": 1}]
    return []


def extract(file_path: str) -> list[dict]:
    """Dispatch to the right extractor based on file extension."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".txt":
        return extract_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def clean(text: str) -> str:
    """Normalize whitespace and strip non-ASCII."""
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^\x00-\x7F]+", " ", text)
    return text.strip()
